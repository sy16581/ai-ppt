from pptx import Presentation
from pptx.util import Inches
from zhipuai import ZhipuAI
from pptx.dml.color import RGBColor

client = ZhipuAI(api_key="622b7be459b291e1fca40f128ae2447c.2tCWCQFAvP5a5ma2")  # 请填写您自己的APIKey

keyword = "ai,应用"
input_content = "ai创新应用"

# 拼接 input 到 title_query
title_query = f"用户想根据{input_content}写一份ppt，你先为他写一个明确的标题，10个字以内，不要给出多个选择，就是一个标题。"

title_response = client.chat.completions.create(
    model="glm-4-flash",  # 请填写您要调用的模型名称
    messages=[
        {"role": "user", "content": title_query},
    ],
)
title = title_response.choices[0].message.content
print(f"title:{title}")


content_query = f'''这是用户的输入：{input_content}
他要就其中的主题写PPT汇报，在PPT之前需要有一份清晰的思路，请你帮他写一份清晰的创作思路，要求尽可能详尽，要求分成四个部分，要返回每部分的思路 和要讲述的内容，下面是示例：
一、引言部分
简单阐述气象学对人类生活的重要性，为整个汇报奠定基调：
提及气象与我们日常生活息息相关，如天气影响着我们的出行、衣着、农业生产、能源消耗等方面。例如，暴雨天气会导致道路积水影响交通，而合适的气象条件对农作物的生长至关重要。
二、气象基础知识部分
对气象学的概念进行简单的解释，让观众对气象学有一个初步的认识：
三、常见气象现象部分
选择一些常见的气象现象进行详细介绍，如降水（雨、雪、冰雹等）、云的类型、雷暴等：
解释每种气象现象的形成机制，包括涉及到的气象要素的变化。
四、气象与人类社会部分
阐述气象对人类社会各个方面的影响，包括正面影响和负面影响。
展望气象学未来的发展趋势，如更精准的预报技术、应对气候变化的策略等。'''

content_response = client.chat.completions.create(
    model="glm-4-flash",  # 请填写您要调用的模型名称
    messages=[
        {"role": "user", "content": content_query},
    ],
)

content = content_response.choices[0].message.content


md_query = f'''根据主题：{keyword} 和标题{title}，参考思路{content}为用户制作一份PPT大纲，以markdown格式输出。大纲必须包括类似示例中对于小标题的内容扩充，不要求详细但需要针对小标题内容简述。
输出格式要求：以markdown格式输出，整个ppt的标题用#，每章的内容用##，每个ppt页的标题用### ，具体内容用-。
输出内容示例：
# 气象知识科普
## 一、引言部分
### 气象学的重要性
- 气象与日常生活的紧密联系
- 天气对出行、衣着、农业、能源的影响
- 实例分析：暴雨与农作物生长
## 二、气象基础知识部分
### 气象学概述
### 气象学的定义
- 气象学的研究领域
- 气象学的基本概念
## 三、常见气象现象部分
### 降水现象
- 降水现象出现的原因
- 降水现象的特点
### 雨的形成
- 雨的形成过程中气象要素的变化
### 雪的形成
- 雪的变化
### 冰雹的形成
- 冰雹的变化
### 云的类型
- 积云的形成与特征
- 层云的形成与特征
- 高层云的形成与特征
### 雷暴
- 雷暴的形成机制
- 雷暴的类型与影响
## 四、气象与人类社会部分
### 气象的正面影响
- 对农业的好处
- 水资源管理的好处
### 气象的负面影响
- 极端天气事件的影响
- 气象灾害的案例
### 气象学的未来发展
- 精准预报技术
### 新技术与发展趋势
- 应对气候变化的策略
气候变化的影响与应对措施'''

# 创建聊天完成的响应，使用指定的模型对用户消息进行回复
md_response = client.chat.completions.create(
    model="glm-4-flash",  # 请填写您要调用的模型名称
    messages=[
        {"role": "user", "content": md_query},
    ],
)

md = md_response.choices[0].message.content

# 读取现有PPT文件
ppt = Presentation('test.pptx')

# 解析Markdown内容，提取标题和对应的内容
section = {'title': None, "child": []}

lines = md.split('\n')
current_main_section = {'title': None, "child": []}
current_sub_section = {'title': None, "child": []}
for line in lines:
    if line.startswith('# '):  # 主标题
        section = {'title': line[2:].strip(), "child": []}
    elif line.startswith('## '):  # 一级标题
        current_section = {'title': line[3:].strip(), "child": []}
        current_main_section = current_section
        section['child'].append(current_section)
    elif line.startswith('### '):  # 二级标题
        current_section = {'title': line[4:].strip(), "child": []}
        current_sub_section = current_section
        current_main_section['child'].append(current_section)
    elif line.startswith('- '):
        current_sub_section['child'].append(line[2:].strip())

def add_text_to_shape(shape, text):
    """将文本添加到形状的占位符中，保留原有样式"""
    """将文本添加到形状的占位符中，保留原有样式"""
    if not shape.has_text_frame:
        return
    text_frame = shape.text_frame

    # 保留原有段落的样式
    if len(text_frame.paragraphs) > 0:
        first_paragraph = text_frame.paragraphs[0]

        # 清除原有文本，但保留样式
        for run in first_paragraph.runs:
            run.text = ""

        # 添加新文本
        run = first_paragraph.add_run()
        run.text = text

        # 复制原有段落的样式
        if len(first_paragraph.runs) > 1:
            original_run = first_paragraph.runs[0]
            run.font.name = original_run.font.name
            run.font.size = original_run.font.size
            run.font.bold = original_run.font.bold
            run.font.italic = original_run.font.italic
            run.font.underline = original_run.font.underline
            # run.font.strike = original_run.font.strike
            # run.font.subscript = original_run.font.subscript
            # run.font.superscript = original_run.font.superscript

    else:
        p = text_frame.add_paragraph()
        p.text = text


def fill_presentation(ppt, section):
    """根据解析后的数据填充PPT"""
    slides = list(ppt.slides)
    index = 0

    def fill_slide(slide, text):
        nonlocal index
        if index < len(slides):
            for shape in slides[index].shapes:
                if shape.has_text_frame:
                    add_text_to_shape(shape, text)
                    break  # 假设每个幻灯片只有一个主要的文本占位符
            index += 1
        else:
            raise ValueError("PPT中的幻灯片数量不足")

    # 填充主标题
    fill_slide(slides[index], section['title'])

    # for main_section in section['child']:
    #     # 填充一级标题
    #     fill_slide(slides[index], main_section['title'])
    #
    #     for sub_section in main_section['child']:
    #         # 填充二级标题
    #         fill_slide(slides[index], sub_section['title'])
    #
    #         for content in sub_section['child']:
    #             # 填充具体内容
    #             fill_slide(slides[index], content)


# 读取现有PPT文件
ppt = Presentation('test.pptx')

# 调用函数填充PPT
fill_presentation(ppt, section)

# 保存PPT文件
ppt.save('updated_test.pptx')
